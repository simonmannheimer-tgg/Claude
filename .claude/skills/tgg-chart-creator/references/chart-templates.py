"""
TGG SEO Deck — Chart Templates
================================
Reusable chart generation code for the TGG Monthly SEO Update deck.
Replace placeholder values in the DATA section of each chart function.
See tgg-chart-design-spec.md for full design rationale and PPTX injection coordinates.

Dependencies: matplotlib, numpy, scipy, pandas, python-pptx
"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.lines import Line2D
from scipy.ndimage import gaussian_filter1d
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Colour tokens ──────────────────────────────────────────────────────────────
BLUE      = '#3c83f6'   # GMC/GSC clicks — matches GSC UI
PURPLE    = '#8762d0'   # GMC/GSC impressions — matches GSC UI
GREEN     = '#1e8c45'   # Shopping / purchases
GREEN_POS = '#70ad47'   # Difference shading: current above prior
RED_NEG   = '#e12b23'   # Difference shading: current below prior
GRAY      = '#7f7f7f'   # Axis labels, tick labels
LGRAY     = '#e8e8e8'   # Grid lines
SPINE_CLR = '#cccccc'   # Axis spines


def apply_base_style(ax):
    """Apply TGG base chart style to an axis."""
    ax.set_facecolor('white')
    ax.tick_params(axis='x', length=0)
    ax.yaxis.grid(True, color=LGRAY, linewidth=0.7, zorder=0)
    ax.xaxis.grid(False)
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(SPINE_CLR)
    ax.spines['bottom'].set_color(SPINE_CLR)


# ══════════════════════════════════════════════════════════════════════════════
# CHART 1 — Multi-month trend with dual axis
# Use for: GMC clicks + impressions over 9 months
# ══════════════════════════════════════════════════════════════════════════════

def build_chart_gmc(output_path: str):
    """
    Dual-axis line chart: clicks (blue, left axis) + impressions (purple, right axis).
    Impressions labels above line, clicks labels below line to prevent collision.
    Smoothed lines, raw values for labels.
    """

    # ── DATA — replace these values each month ─────────────────────────────
    months     = ['Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
    clicks_raw = [263240, 277634, 276849, 330040, 469548, 524646, 363092, 265120, 265111]
    impr_raw   = [14069035, 15062816, 14564541, 16710454, 23401679,
                  23985665, 19542189, 14095742, 13999638]
    # ── END DATA ────────────────────────────────────────────────────────────

    SIG = 1.2  # gaussian smoothing sigma — adjust if curve feels too flat or jagged
    x   = np.arange(len(months))
    clicks_s = gaussian_filter1d(np.array(clicks_raw, dtype=float), sigma=SIG)
    impr_s   = gaussian_filter1d(np.array(impr_raw,   dtype=float), sigma=SIG)

    fig, ax1 = plt.subplots(figsize=(6.452, 2.421))
    fig.patch.set_facecolor('white')
    apply_base_style(ax1)
    ax1.set_xlim(-0.5, len(months) - 0.5)
    ax1.set_ylim(150, 720)

    # Clicks line (left axis)
    ax1.plot(x, clicks_s / 1000, color=BLUE, linewidth=2.0, zorder=4, label='Clicks')
    ax1.set_ylabel('Clicks', fontsize=7.5, color=BLUE)
    ax1.tick_params(axis='y', colors=BLUE, labelsize=7)
    ax1.set_xticks(x)
    ax1.set_xticklabels(months, fontsize=8, color=GRAY)

    # Impressions line (right axis)
    ax2 = ax1.twinx()
    ax2.set_facecolor('none')
    ax2.plot(x, impr_s / 1e6, color=PURPLE, linewidth=1.8, zorder=3, label='Impressions')
    ax2.set_ylabel('Impressions (M)', fontsize=7.5, color=PURPLE)
    ax2.tick_params(axis='y', colors=PURPLE, labelsize=7)
    ax2.set_ylim(5, 38)
    for sp in ['top', 'left', 'bottom']:
        ax2.spines[sp].set_visible(False)
    ax2.spines['right'].set_color(SPINE_CLR)

    n = len(months)

    # Impressions labels ABOVE the impressions line — raw values
    for i, (xi, iv) in enumerate(zip(x, impr_raw)):
        ax2.annotate(
            f'{iv / 1e6:.1f}M',
            xy=(xi, impr_s[i] / 1e6),
            xytext=(0, 7), textcoords='offset points',
            ha='center', va='bottom', fontsize=6.0, color=PURPLE,
            fontweight='bold' if i == n - 1 else 'normal'
        )

    # Clicks labels BELOW the clicks line — raw values
    for i, (xi, cv) in enumerate(zip(x, clicks_raw)):
        ax2.annotate(
            f'{cv / 1000:.0f}K',
            xy=(xi, impr_s[i] / 1e6),
            xytext=(0, -22), textcoords='offset points',
            ha='center', va='top', fontsize=6.0, color=BLUE,
            fontweight='bold' if i == n - 1 else 'normal'
        )

    # Legend
    h1 = Line2D([0], [0], color=BLUE,   linewidth=2.0, label='Clicks')
    h2 = Line2D([0], [0], color=PURPLE, linewidth=1.8, label='Impressions')
    ax1.legend(handles=[h1, h2], fontsize=7.5, frameon=False,
               loc='upper left', labelcolor=GRAY)

    plt.tight_layout(pad=0.3)
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


# ══════════════════════════════════════════════════════════════════════════════
# CHART 2 — Current vs prior month day-by-day with difference shading
# Use for: organic shopping purchases, revenue, CVR — any daily metric MoM
# ══════════════════════════════════════════════════════════════════════════════

def build_chart_purchases(output_path: str):
    """
    Two-line chart: current month (solid) vs prior month (dashed).
    Difference shading: green where current > prior, red where current < prior.
    Both lines in GREEN. Normalised to shorter month's day count.
    """

    # ── DATA — replace these values each month ─────────────────────────────
    # Daily purchases — list one value per day, sorted day 1 to day N
    # Feb 2026 (28 days) and Mar 2026 (31 days) — truncated to 28 for comparison
    prior_daily = [
        214, 157, 184, 176, 159, 136, 145, 160, 125, 131,
        144, 158, 180, 139, 122, 144, 151, 166, 189, 136,
        184, 165, 182, 141, 125, 178, 125, 173
    ]  # February 2026 — 28 values
    current_daily = [
        153, 110, 130, 122, 122,  99,  95, 120, 132, 109,
        107, 126, 111, 122, 151, 156, 135, 143, 196, 159,
        137, 215, 170, 137, 152, 149, 143, 142
    ]  # March 2026 — first 28 days
    prior_label   = 'February'
    current_label = 'March'
    ylabel        = 'Purchases'
    # ── END DATA ────────────────────────────────────────────────────────────

    SIG = 1.4
    n   = min(len(prior_daily), len(current_daily))
    x   = np.arange(1, n + 1)
    prior_s   = gaussian_filter1d(np.array(prior_daily[:n],   dtype=float), sigma=SIG)
    current_s = gaussian_filter1d(np.array(current_daily[:n], dtype=float), sigma=SIG)

    fig, ax = plt.subplots(figsize=(6.452, 2.421))
    fig.patch.set_facecolor('white')
    apply_base_style(ax)
    ax.set_xlim(1, n + 2.5)
    ax.set_xticks([1, 7, 14, 21, n])

    # Difference shading
    ax.fill_between(x, prior_s, current_s,
                    where=(current_s >= prior_s), interpolate=True,
                    color=GREEN_POS, alpha=0.15, zorder=2)
    ax.fill_between(x, prior_s, current_s,
                    where=(current_s < prior_s), interpolate=True,
                    color=RED_NEG, alpha=0.12, zorder=2)

    # Lines
    ax.plot(x, prior_s,   color=GREEN, linewidth=1.4, linestyle='--',
            alpha=0.5, zorder=3, label=prior_label)
    ax.plot(x, current_s, color=GREEN, linewidth=2.0,
            zorder=4, label=current_label)

    # End labels
    ax.text(n + 0.4, prior_s[-1],   prior_label,   fontsize=7,   color=GREEN, alpha=0.6, va='center')
    ax.text(n + 0.4, current_s[-1], current_label, fontsize=7.5, color=GREEN, fontweight='bold', va='center')

    ax.set_ylabel(ylabel, fontsize=7.5, color=GRAY)
    ax.set_xlabel('Day of month', fontsize=7.5, color=GRAY)
    ax.tick_params(labelsize=7, colors=GRAY)
    ax.legend(fontsize=7.5, frameon=False, labelcolor=GRAY, loc='upper right')

    plt.tight_layout(pad=0.3)
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


# ══════════════════════════════════════════════════════════════════════════════
# PPTX INJECTION — Slide 3
# ══════════════════════════════════════════════════════════════════════════════

def inject_charts_into_deck(
    pptx_path: str,
    output_path: str,
    chart_top_path: str,
    chart_bot_path: str,
    header_top: str = 'GMC Clicks & Impressions  Jul 2025 to Mar 2026',
    header_bot: str = 'Organic Shopping Purchases  Feb vs Mar 2026',
    slide_index: int = 2,          # 0-indexed; slide 3 = index 2
    remove_bullet_text: str = None # partial string match to remove a bullet
):
    """
    Removes IMAGE placeholder shapes from the target slide, adds chart images
    and text headers. Optionally removes a bullet point by partial text match.

    Placeholder positions (EMU) — derived from original deck XML:
        Top slot:    left=824555  top=1920575  w=5899428  h=2214064
        Bottom slot: left=824555  top=4362198  w=5899428  h=2214064
    """
    HEADER_H = Emu(228600)   # 0.25" header height
    GAP      = Emu(57150)    # 0.0625" gap between header and chart

    prs   = Presentation(pptx_path)
    slide = prs.slides[slide_index]

    # Remove specified bullet
    if remove_bullet_text:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if remove_bullet_text in ''.join(r.text for r in para.runs):
                    para._p.getparent().remove(para._p)
                    break

    # Find and remove IMAGE placeholders
    image_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if any('IMAGE' in r.text for r in para.runs):
                    image_shapes.append(shape)
                    break
    image_shapes.sort(key=lambda s: s.top)
    for s in image_shapes:
        s._element.getparent().remove(s._element)

    slots = [
        (image_shapes[0].left, image_shapes[0].top,
         image_shapes[0].width, image_shapes[0].height, header_top, chart_top_path),
        (image_shapes[1].left, image_shapes[1].top,
         image_shapes[1].width, image_shapes[1].height, header_bot, chart_bot_path),
    ]

    for left, top, width, height, header_text, chart_path in slots:
        # Header text box
        txBox = slide.shapes.add_textbox(left, top, width, HEADER_H)
        tf    = txBox.text_frame
        tf.word_wrap = False
        p    = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run  = p.add_run()
        run.text = header_text
        run.font.size  = Pt(9)
        run.font.bold  = True
        run.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
        run.font.name  = 'Calibri'

        # Chart image
        chart_top  = top + HEADER_H + GAP
        chart_h    = height - HEADER_H - GAP
        slide.shapes.add_picture(chart_path, left, chart_top, width, chart_h)

    prs.save(output_path)
    print(f"Saved: {output_path}")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN — run to regenerate both charts and inject into deck
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    GMC_CHART_PATH       = '/tmp/chart_gmc.png'
    PURCHASES_CHART_PATH = '/tmp/chart_purchases.png'
    INPUT_DECK           = 'TGG_MARCH_Close_to_finished.pptx'
    OUTPUT_DECK          = 'TGG_March_2026_final.pptx'

    print("Building GMC chart...")
    build_chart_gmc(GMC_CHART_PATH)

    print("Building purchases chart...")
    build_chart_purchases(PURCHASES_CHART_PATH)

    print("Injecting into deck...")
    inject_charts_into_deck(
        pptx_path       = INPUT_DECK,
        output_path     = OUTPUT_DECK,
        chart_top_path  = GMC_CHART_PATH,
        chart_bot_path  = PURCHASES_CHART_PATH,
        remove_bullet_text = 'JB keyword footprint'
    )
    print("Done.")
